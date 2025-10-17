from pptx import Presentation
from pptx.util import Inches, Pt
import os
import logging
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN


# Logging sozlamalari
logger = logging.getLogger(__name__)

def replace_text_and_keep_format(slide, old_text: str, new_text: str):
       
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        cleaned_new_text = new_text.strip().replace('\n', ' ')

        for paragraph in shape.text_frame.paragraphs:
            if old_text in shape.text_frame.text:
                if len(paragraph.runs) == 1:
                    run = paragraph.runs[0]
                    font_name = run.font.name
                    font_size = run.font.size
                    font_bold = run.font.bold
                    paragraph_alignment = paragraph.alignment 
                    run.text = run.text.replace(old_text, cleaned_new_text)
                    
                    if paragraph.runs:
                        new_run = paragraph.runs[0]
                        if font_name:
                            new_run.font.name = font_name
                        if font_size:
                            new_run.font.size = font_size
                        if font_bold is not None:
                            new_run.font.bold = font_bold

                    if paragraph_alignment:
                        paragraph.alignment = paragraph_alignment
                else:
                    paragraph_alignment = paragraph.alignment
                    paragraph.text = paragraph.text.replace(old_text, cleaned_new_text)

                    if paragraph_alignment:
                        paragraph.alignment = paragraph_alignment    

def sertifikat_yaratish(
    shablon_fayli: str,
    chiqish_fayli: str,
    student_ismi: str,
    fan: str,
    foiz: str,
    teacher: str,
    orin: str
):

    
    if not os.path.exists(shablon_fayli):
        logger.error(f"Shablon fayli topilmadi: {shablon_fayli}")
        raise FileNotFoundError(f"Shablon fayli topilmadi: {shablon_fayli}")

    try:
        # Shablon faylni ochish
        prs = Presentation(shablon_fayli)
        replacements = {
            "STUDENT_ISM": student_ismi,  # Talaba F.I.SH. uchun joy tutuvchi
            "(fan)": fan,           # Fan nomi uchun joy tutuvchi
            "(foiz)": foiz,         # Natija foizi uchun joy tutuvchi
            "(teacher)": teacher,   # O'qituvchi F.I.SH. uchun joy tutuvchi
            "(orin)": orin,         # Egallangan o'rin uchun joy tutuvchi
        }

        # Slaydlarni aylanib chiqish
        for slide in prs.slides:
            for key, value in replacements.items():
                replace_text_and_keep_format(slide, key, value)


        prs.save(chiqish_fayli)
        logger.info(f"Sertifikat muvaffaqiyatli yaratildi: {chiqish_fayli}")
        return True

    except Exception as e:
        logger.error(f"Sertifikat yaratishda xato: {e}")
        raise e

# Agar fayl bevosita ishga tushirilsa (sinov uchun)
if __name__ == '__main__':
    # Bu qism faqat funksiyani sinash uchun ishlatiladi.
    # Loyiha ishga tushirilganda bu qism ishlamaydi.
    print("Sertifikat generatori funksiyasi tayyor. Iltimos, ishlatishdan oldin 'sertifikat1.pptx' shablonini yarating.")
